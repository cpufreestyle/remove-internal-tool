import os
import sys
import csv as csvlib
import warnings
import contextlib
import threading
import tkinter as tk
from tkinter import ttk, filedialog, messagebox
from pathlib import Path
import string
from ctypes import windll

warnings.filterwarnings("ignore")
try:
    from PyPDF2.errors import PdfReadWarning
    warnings.filterwarnings("ignore", category=PdfReadWarning)
except Exception:
    pass

from docx import Document
from openpyxl import load_workbook
from PyPDF2 import PdfReader, PdfWriter
from pptx import Presentation

DOC_TYPES = ['.docx', '.xlsx', '.pdf', '.ppt', '.pptx', '.txt', '.csv', '.md', '.rtf']
IMG_TYPES = ['.jpg', '.jpeg', '.png', '.bmp', '.tiff', '.tif', '.webp']
FILE_TYPES = DOC_TYPES + IMG_TYPES

ocr_engine = None

def init_ocr():
    global ocr_engine
    if ocr_engine is not None:
        return True
    try:
        from paddleocr import PaddleOCR
        ocr_engine = PaddleOCR(use_angle_cls=True, lang='ch', show_log=False)
        return True
    except Exception:
        return False

def get_all_drives():
    drives = []
    bitmask = windll.kernel32.GetLogicalDrives()
    for letter in string.ascii_uppercase:
        if bitmask & 1:
            drive = f'{letter}:\\'
            try:
                import win32file
                drive_type = win32file.GetDriveType(drive)
                if drive_type in [2, 5]:
                    bitmask >>= 1
                    continue
            except Exception:
                pass
            drives.append(drive)
        bitmask >>= 1
    return drives

def process_docx(file_path, changed_files):
    doc = Document(file_path)
    changed = False
    if doc.paragraphs:
        new_text = doc.paragraphs[0].text.replace('内部', '')
        if new_text != doc.paragraphs[0].text:
            doc.paragraphs[0].text = new_text
            changed = True
    doc.save(file_path)
    if changed:
        changed_files.append(file_path)

def process_xlsx(file_path, changed_files):
    wb = load_workbook(file_path)
    ws = wb.active
    first_row = next(ws.iter_rows(min_row=1, max_row=1))
    changed = False
    for cell in first_row:
        if cell.value and isinstance(cell.value, str):
            new_val = cell.value.replace('内部', '')
            if new_val != cell.value:
                cell.value = new_val
                changed = True
    wb.save(file_path)
    if changed:
        changed_files.append(file_path)

def process_pdf(file_path, changed_files):
    with open(os.devnull, 'w') as devnull, contextlib.redirect_stderr(devnull):
        reader = PdfReader(file_path)
        writer = PdfWriter()
        changed = False
        if reader.pages:
            first_page = reader.pages[0]
            text = first_page.extract_text()
            if text:
                lines = text.split('\n')
                if lines and '内部' in lines[0]:
                    changed = True
            writer.add_page(first_page)
            for page in reader.pages[1:]:
                writer.add_page(page)
            with open(file_path, 'wb') as f:
                writer.write(f)
        if changed:
            changed_files.append(file_path)

def process_ppt(file_path, changed_files):
    prs = Presentation(file_path)
    changed = False
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text:
                lines = shape.text.split('\n')
                if lines and '内部' in lines[0]:
                    lines[0] = lines[0].replace('内部', '')
                    shape.text = '\n'.join(lines)
                    changed = True
    prs.save(file_path)
    if changed:
        changed_files.append(file_path)

def process_txt(file_path, changed_files):
    for enc in ['utf-8', 'gbk', 'gb2312']:
        try:
            with open(file_path, 'r', encoding=enc) as f:
                lines = f.readlines()
            if lines and '内部' in lines[0]:
                lines[0] = lines[0].replace('内部', '')
                with open(file_path, 'w', encoding=enc) as f:
                    f.writelines(lines)
                changed_files.append(file_path)
            return
        except Exception:
            continue

def process_image(file_path, changed_files):
    from PIL import Image
    img = Image.open(file_path)
    if not init_ocr():
        return
    result = ocr_engine.ocr(file_path, cls=True)
    if not result or not result[0]:
        return
    changed = False
    draw = None
    for line in result[0]:
        box = line[0]
        text = line[1][0]
        if '内部' in text:
            if draw is None:
                from PIL import ImageDraw
                draw = ImageDraw.Draw(img)
            x1, y1 = box[0]
            x2, y2 = box[2]
            draw.rectangle([x1, y1, x2, y2], fill='white')
            changed = True
    if changed:
        img.save(file_path)
        changed_files.append(file_path)


class App(tk.Tk):
    def __init__(self):
        super().__init__()
        self.title("全盘文档批量去除[内部]工具 v3.0")
        self.geometry("720x560")
        self.resizable(True, True)
        self.configure(bg="#1e1e2e")
        self._running = False
        self._stop_flag = False
        self._build_ui()

    def _build_ui(self):
        # ── 顶部标题 ──
        header = tk.Frame(self, bg="#313244", pady=12)
        header.pack(fill="x")
        tk.Label(header, text="🗂  全盘文档批量去除[内部]工具",
                 font=("微软雅黑", 16, "bold"),
                 fg="#cdd6f4", bg="#313244").pack()
        tk.Label(header, text="支持 docx / xlsx / pdf / ppt / txt / csv / md / rtf / 图片",
                 font=("微软雅黑", 9), fg="#a6adc8", bg="#313244").pack()

        # ── 扫描目录选择 ──
        dir_frame = tk.Frame(self, bg="#1e1e2e", pady=8)
        dir_frame.pack(fill="x", padx=20)
        tk.Label(dir_frame, text="扫描目录：", font=("微软雅黑", 10),
                 fg="#cdd6f4", bg="#1e1e2e").pack(side="left")
        self.dir_var = tk.StringVar(value="全盘扫描（所有磁盘）")
        self.dir_entry = tk.Entry(dir_frame, textvariable=self.dir_var,
                                  font=("微软雅黑", 10), width=38,
                                  bg="#313244", fg="#cdd6f4",
                                  insertbackground="#cdd6f4",
                                  relief="flat", bd=4)
        self.dir_entry.pack(side="left", padx=(4, 6))
        tk.Button(dir_frame, text="选择目录", font=("微软雅黑", 9),
                  bg="#89b4fa", fg="#1e1e2e", relief="flat",
                  padx=8, command=self._choose_dir).pack(side="left", padx=2)
        tk.Button(dir_frame, text="全盘", font=("微软雅黑", 9),
                  bg="#a6e3a1", fg="#1e1e2e", relief="flat",
                  padx=8, command=self._reset_dir).pack(side="left", padx=2)

        # ── 统计栏 ──
        stats_frame = tk.Frame(self, bg="#1e1e2e")
        stats_frame.pack(fill="x", padx=20, pady=(0, 6))
        self.stat_scanned = self._stat_label(stats_frame, "已扫描", "0")
        self.stat_changed = self._stat_label(stats_frame, "已修改", "0")
        self.stat_status  = self._stat_label(stats_frame, "状态", "就绪")

        # ── 进度条 ──
        self.progress = ttk.Progressbar(self, mode="indeterminate",
                                        style="green.Horizontal.TProgressbar")
        self.progress.pack(fill="x", padx=20, pady=(0, 6))
        style = ttk.Style(self)
        style.theme_use("default")
        style.configure("green.Horizontal.TProgressbar",
                        troughcolor="#313244", background="#a6e3a1",
                        thickness=8)

        # ── 日志框 ──
        log_frame = tk.Frame(self, bg="#1e1e2e")
        log_frame.pack(fill="both", expand=True, padx=20, pady=(0, 8))
        self.log_box = tk.Text(log_frame, font=("Consolas", 9),
                               bg="#181825", fg="#cdd6f4",
                               insertbackground="#cdd6f4",
                               relief="flat", bd=0,
                               state="disabled", wrap="none")
        sb_y = tk.Scrollbar(log_frame, command=self.log_box.yview)
        sb_x = tk.Scrollbar(log_frame, orient="horizontal",
                             command=self.log_box.xview)
        self.log_box.configure(yscrollcommand=sb_y.set,
                               xscrollcommand=sb_x.set)
        sb_y.pack(side="right", fill="y")
        sb_x.pack(side="bottom", fill="x")
        self.log_box.pack(fill="both", expand=True)

        # ── 按钮区 ──
        btn_frame = tk.Frame(self, bg="#1e1e2e", pady=8)
        btn_frame.pack()
        self.btn_start = tk.Button(btn_frame, text="▶  开始扫描",
                                   font=("微软雅黑", 11, "bold"),
                                   bg="#89b4fa", fg="#1e1e2e",
                                   relief="flat", padx=20, pady=6,
                                   command=self._start)
        self.btn_start.pack(side="left", padx=8)
        self.btn_stop = tk.Button(btn_frame, text="⏹  停止",
                                  font=("微软雅黑", 11),
                                  bg="#f38ba8", fg="#1e1e2e",
                                  relief="flat", padx=20, pady=6,
                                  state="disabled",
                                  command=self._stop)
        self.btn_stop.pack(side="left", padx=8)
        self.btn_export = tk.Button(btn_frame, text="📄  导出报告",
                                    font=("微软雅黑", 11),
                                    bg="#a6e3a1", fg="#1e1e2e",
                                    relief="flat", padx=20, pady=6,
                                    state="disabled",
                                    command=self._export)
        self.btn_export.pack(side="left", padx=8)

        self._changed_files = []
        self._scanned = 0

    def _stat_label(self, parent, title, value):
        f = tk.Frame(parent, bg="#313244", padx=14, pady=6)
        f.pack(side="left", padx=6, pady=4)
        tk.Label(f, text=title, font=("微软雅黑", 8),
                 fg="#a6adc8", bg="#313244").pack()
        lbl = tk.Label(f, text=value, font=("微软雅黑", 14, "bold"),
                       fg="#cdd6f4", bg="#313244")
        lbl.pack()
        return lbl

    def _choose_dir(self):
        d = filedialog.askdirectory(title="选择扫描目录")
        if d:
            self.dir_var.set(d)

    def _reset_dir(self):
        self.dir_var.set("全盘扫描（所有磁盘）")

    def _log(self, msg, color="#cdd6f4"):
        self.log_box.configure(state="normal")
        self.log_box.insert("end", msg + "\n")
        self.log_box.see("end")
        self.log_box.configure(state="disabled")

    def _start(self):
        if self._running:
            return
        self._running = True
        self._stop_flag = False
        self._changed_files = []
        self._scanned = 0
        self.log_box.configure(state="normal")
        self.log_box.delete("1.0", "end")
        self.log_box.configure(state="disabled")
        self.stat_scanned.config(text="0")
        self.stat_changed.config(text="0")
        self.stat_status.config(text="扫描中…")
        self.btn_start.config(state="disabled")
        self.btn_stop.config(state="normal")
        self.btn_export.config(state="disabled")
        self.progress.start(12)
        threading.Thread(target=self._run, daemon=True).start()

    def _stop(self):
        self._stop_flag = True
        self._log("⏹ 用户请求停止，等待当前文件处理完成…", "#fab387")

    def _run(self):
        dir_val = self.dir_var.get()
        if dir_val == "全盘扫描（所有磁盘）":
            roots = get_all_drives()
        else:
            roots = [dir_val]

        self._log(f"🚀 开始扫描：{', '.join(roots)}")
        SKIP = ["onedrive", "cloud", "dropbox", "baidunetdisk", "坚果云",
                "sync", "googledrive", "$recycle.bin",
                "system volume information", "windows", "program files"]

        for root in roots:
            if self._stop_flag:
                break
            self._log(f"📂 磁盘：{root}")
            for dirpath, _, filenames in os.walk(root):
                if self._stop_flag:
                    break
                dl = dirpath.lower()
                if any(x in dl for x in SKIP):
                    continue
                for filename in filenames:
                    if self._stop_flag:
                        break
                    ext = os.path.splitext(filename)[1].lower()
                    if ext not in FILE_TYPES:
                        continue
                    fp = os.path.join(dirpath, filename)
                    try:
                        before = len(self._changed_files)
                        if ext == '.docx':
                            process_docx(fp, self._changed_files)
                        elif ext == '.xlsx':
                            process_xlsx(fp, self._changed_files)
                        elif ext == '.pdf':
                            process_pdf(fp, self._changed_files)
                        elif ext in ('.ppt', '.pptx'):
                            process_ppt(fp, self._changed_files)
                        elif ext in ('.txt', '.csv', '.md', '.rtf'):
                            process_txt(fp, self._changed_files)
                        elif ext in IMG_TYPES:
                            process_image(fp, self._changed_files)
                        self._scanned += 1
                        if len(self._changed_files) > before:
                            self._log(f"  ✅ 已修改：{fp}", "#a6e3a1")
                        if self._scanned % 50 == 0:
                            self.after(0, self._update_stats)
                    except Exception as e:
                        self._log(f"  ⚠ 跳过：{fp}  ({e})", "#f38ba8")

        self.after(0, self._done)

    def _update_stats(self):
        self.stat_scanned.config(text=str(self._scanned))
        self.stat_changed.config(text=str(len(self._changed_files)))

    def _done(self):
        self._running = False
        self.progress.stop()
        self._update_stats()
        status = "已停止" if self._stop_flag else "完成 ✓"
        self.stat_status.config(text=status)
        self.btn_start.config(state="normal")
        self.btn_stop.config(state="disabled")
        self.btn_export.config(state="normal")
        self._log(f"\n{'='*50}")
        self._log(f"🎉 扫描{status}！共扫描 {self._scanned} 个文件，修改 {len(self._changed_files)} 个")
        self._log(f"{'='*50}")

    def _export(self):
        if not self._changed_files:
            messagebox.showinfo("提示", "没有修改过的文件，无需导出。")
            return
        save_path = filedialog.asksaveasfilename(
            defaultextension=".csv",
            filetypes=[("CSV 文件", "*.csv")],
            initialfile="changed_files.csv",
            title="保存报告"
        )
        if not save_path:
            return
        with open(save_path, "w", newline='', encoding="utf-8-sig") as f:
            w = csvlib.writer(f)
            w.writerow(["文件路径"])
            for fp in self._changed_files:
                w.writerow([fp])
        messagebox.showinfo("导出成功", f"报告已保存：\n{save_path}")
        self._log(f"📄 报告已导出：{save_path}", "#89b4fa")


if __name__ == "__main__":
    app = App()
    app.mainloop()
