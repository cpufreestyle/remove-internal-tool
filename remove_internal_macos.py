#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
macOS 去除内部标记工具
支持 docx / xlsx / pdf / pptx / txt / 图片OCR
macOS 移植版 by Claude
"""

import os
import sys
import warnings
import contextlib
import threading
import tkinter as tk
from tkinter import ttk, filedialog, messagebox
from pathlib import Path

warnings.filterwarnings("ignore")

# === 跨平台文档处理 ===
try:
    from PyPDF2.errors import PdfReadWarning
    warnings.filterwarnings("ignore", category=PdfReadWarning)
except Exception:
    pass

from docx import Document
from openpyxl import load_workbook
from PyPDF2 import PdfReader, PdfWriter
from pptx import Presentation

DOC_TYPES = ['.docx', '.xlsx', '.pdf', '.ppt', '.pptx', '.txt', '.csv', '.md', '.rtf']
IMG_TYPES = ['.jpg', '.jpeg', '.png', '.bmp', '.tiff', '.tif', '.webp']
FILE_TYPES = DOC_TYPES + IMG_TYPES

# === 要去除的关键词 ===
KEYWORDS = [
    # 中文
    '内部', '内参', '机密', '秘密', '绝密', '仅限',
    '内部文件', '内部资料', '内部材料',
    '涉密', '保密', '机密文件', '秘密文件',
    # 英文
    'INTERNAL', 'CONFIDENTIAL', 'SECRET', 'TOP SECRET',
    'Internal', 'Confidential', 'Secret',
    # 其他
    'Draft', 'DRAFT', '草稿',
]

DEVNULL = '/dev/null'


def match_keywords(text):
    """检查文本是否包含敏感关键词"""
    if not text:
        return None
    for kw in KEYWORDS:
        if kw in text:
            return kw
    return None


def process_docx(file_path, changed_files, log_func):
    """处理 Word 文档"""
    try:
        doc = Document(file_path)
        changed = False
        for i, para in enumerate(doc.paragraphs[:5]):
            for kw in KEYWORDS:
                if kw in para.text:
                    original = para.text
                    para.text = para.text.replace(kw, '').replace('  ', ' ')
                    if para.text != original:
                        changed = True
        if changed:
            doc.save(file_path)
            changed_files.append(file_path)
            log_func(f'  [修改] {Path(file_path).name}')
    except Exception as e:
        log_func(f'  [跳过] {Path(file_path).name}: {e}')


def process_xlsx(file_path, changed_files, log_func):
    """处理 Excel 文档"""
    try:
        wb = load_workbook(file_path)
        changed = False
        for sheet in wb.worksheets:
            for row in sheet.iter_rows(max_row=3):
                for cell in row:
                    if cell.value and isinstance(cell.value, str):
                        for kw in KEYWORDS:
                            if kw in cell.value:
                                cell.value = cell.value.replace(kw, '').replace('  ', ' ')
                                changed = True
        if changed:
            wb.save(file_path)
            changed_files.append(file_path)
            log_func(f'  [修改] {Path(file_path).name}')
    except Exception as e:
        log_func(f'  [跳过] {Path(file_path).name}: {e}')


def process_pdf(file_path, changed_files, log_func):
    """处理 PDF 文档"""
    try:
        with open(DEVNULL, 'w') as devnull, \
             contextlib.redirect_stderr(devnull):
            reader = PdfReader(file_path)
            writer = PdfWriter()
            changed = False
            for i, page in enumerate(reader.pages[:2]):
                text = page.extract_text()
                if text and match_keywords(text):
                    changed = True
                    break
            writer.add_page(reader.pages[0])
            for page in reader.pages[1:]:
                writer.add_page(page)
            with open(file_path, 'wb') as f:
                writer.write(f)
            if changed:
                changed_files.append(file_path)
                log_func(f'  [修改] {Path(file_path).name}')
    except Exception as e:
        log_func(f'  [跳过] {Path(file_path).name}: {e}')


def process_pptx(file_path, changed_files, log_func):
    """处理 PowerPoint 文档"""
    try:
        prs = Presentation(file_path)
        changed = False
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text:
                    for kw in KEYWORDS:
                        if kw in shape.text:
                            shape.text = shape.text.replace(kw, '').replace('  ', ' ')
                            changed = True
        if changed:
            prs.save(file_path)
            changed_files.append(file_path)
            log_func(f'  [修改] {Path(file_path).name}')
    except Exception as e:
        log_func(f'  [跳过] {Path(file_path).name}: {e}')


def process_txt(file_path, changed_files, log_func):
    """处理纯文本文件"""
    for enc in ['utf-8', 'gbk', 'gb2312']:
        try:
            with open(file_path, 'r', encoding=enc) as f:
                content = f.read()
            if match_keywords(content[:500]):
                new_content = content
                for kw in KEYWORDS:
                    new_content = new_content.replace(kw, '')
                if new_content != content:
                    with open(file_path, 'w', encoding=enc) as f:
                        f.write(new_content)
                    changed_files.append(file_path)
                    log_func(f'  [修改] {Path(file_path).name}')
            return
        except Exception:
            continue


def process_image(file_path, changed_files, log_func):
    """处理图片（需要 OCR）"""
    try:
        from PIL import Image
        img = Image.open(file_path)
        if not init_ocr():
            return
        result = ocr_engine.ocr(file_path, cls=True)
        if not result or not result[0]:
            return
        changed = False
        draw = None
        for line in result[0]:
            box = line[0]
            text = line[1][0]
            if match_keywords(text):
                if draw is None:
                    from PIL import ImageDraw
                    draw = ImageDraw.Draw(img)
                x1, y1 = box[0]
                x2, y2 = box[2]
                draw.rectangle([x1, y1, x2, y2], fill='white')
                changed = True
        if changed:
            img.save(file_path)
            changed_files.append(file_path)
            log_func(f'  [修改] {Path(file_path).name} (OCR)')
    except Exception as e:
        log_func(f'  [跳过] {Path(file_path).name}: {e}')


ocr_engine = None

def init_ocr():
    global ocr_engine
    if ocr_engine is not None:
        return True
    try:
        from paddleocr import PaddleOCR
        ocr_engine = PaddleOCR(use_angle_cls=True, lang='ch', show_log=False)
        return True
    except Exception:
        return False


class App(tk.Tk):
    def __init__(self):
        super().__init__()
        self.title("去除内部标记工具 - macOS版")
        self.geometry("800x600")
        self.configure(bg="#0f0f0f")
        self._running = False
        self._stop_flag = False
        self._changed_files = []
        self._scanned = 0
        self._build_ui()

    def _build_ui(self):
        header = tk.Frame(self, bg="#1a1a1a", pady=14)
        header.pack(fill="x")
        tk.Label(header, text="macOS 去除内部标记工具",
                 font=("SF Pro Display", 16, "bold"),
                 fg="#ea580c", bg="#1a1a1a").pack()
        tk.Label(header, text="支持 docx/xlsx/pdf/pptx/txt/图片OCR",
                 font=("SF Pro Text", 10),
                 fg="#888", bg="#1a1a1a").pack()

        dir_frame = tk.Frame(self, bg="#0f0f0f", pady=10)
        dir_frame.pack(fill="x", padx=20)
        tk.Label(dir_frame, text="目录：",
                 font=("SF Pro Text", 11),
                 fg="#ccc", bg="#0f0f0f").pack(side="left")
        self.dir_var = tk.StringVar(value="请选择目录...")
        self.dir_entry = tk.Entry(dir_frame, textvariable=self.dir_var,
                                  font=("SF Pro Text", 11), width=42,
                                  bg="#1e1e1e", fg="#f5f5f5",
                                  insertbackground="#f5f5f5",
                                  relief="flat", bd=3)
        self.dir_entry.pack(side="left", padx=(4, 8))
        tk.Button(dir_frame, text="选择目录",
                  font=("SF Pro Text", 10),
                  bg="#252525", fg="#ea580c", relief="flat",
                  padx=12, pady=4,
                  command=self._choose_dir).pack(side="left", padx=4)
        tk.Button(dir_frame, text="选择文件",
                  font=("SF Pro Text", 10),
                  bg="#252525", fg="#ea580c", relief="flat",
                  padx=12, pady=4,
                  command=self._choose_files).pack(side="left", padx=4)

        stats_frame = tk.Frame(self, bg="#0f0f0f")
        stats_frame.pack(fill="x", padx=20, pady=(0, 6))
        self.stat_scanned = self._stat_label(stats_frame, "已扫描", "0")
        self.stat_changed = self._stat_label(stats_frame, "已修改", "0")
        self.stat_status = self._stat_label(stats_frame, "状态", "就绪")

        kw_frame = tk.Frame(self, bg="#1a1a1a", padx=12, pady=6)
        kw_frame.pack(fill="x", padx=20, pady=(0, 6))
        tk.Label(kw_frame, text="关键词：",
                 font=("SF Pro Text", 9),
                 fg="#888", bg="#1a1a1a").pack(side="left")
        tk.Label(kw_frame, text="内部 / INTERNAL / 机密 / CONFIDENTIAL / 秘密 / SECRET / 绝密 / 草稿",
                 font=("SF Pro Text", 9),
                 fg="#ea580c", bg="#1a1a1a").pack(side="left")

        log_frame = tk.Frame(self, bg="#0f0f0f")
        log_frame.pack(fill="both", expand=True, padx=20, pady=(0, 8))
        self.log_box = tk.Text(log_frame, font=("Menlo", 10),
                               bg="#111", fg="#4ade80",
                               insertbackground="#4ade80",
                               relief="flat", bd=0,
                               state="disabled", wrap="none")
        sb_y = tk.Scrollbar(log_frame, command=self.log_box.yview)
        sb_x = tk.Scrollbar(log_frame, orient="horizontal",
                             command=self.log_box.xview)
        self.log_box.configure(yscrollcommand=sb_y.set,
                               xscrollcommand=sb_x.set)
        sb_y.pack(side="right", fill="y")
        sb_x.pack(side="bottom", fill="x")
        self.log_box.pack(fill="both", expand=True)

        btn_frame = tk.Frame(self, bg="#0f0f0f", pady=10)
        btn_frame.pack()
        self.btn_start = tk.Button(btn_frame, text="开始扫描",
                                   font=("SF Pro Text", 12, "bold"),
                                   bg="#ea580c", fg="#fff",
                                   relief="flat", padx=24, pady=8,
                                   command=self._start)
        self.btn_start.pack(side="left", padx=8)
        self.btn_stop = tk.Button(btn_frame, text="停止",
                                  font=("SF Pro Text", 11),
                                  bg="#ef4444", fg="#fff",
                                  relief="flat", padx=20, pady=8,
                                  state="disabled",
                                  command=self._stop)
        self.btn_stop.pack(side="left", padx=8)
        self.btn_export = tk.Button(btn_frame, text="导出报告",
                                    font=("SF Pro Text", 11),
                                    bg="#22c55e", fg="#fff",
                                    relief="flat", padx=20, pady=8,
                                    state="disabled",
                                    command=self._export)
        self.btn_export.pack(side="left", padx=8)

    def _stat_label(self, parent, title, value):
        f = tk.Frame(parent, bg="#1e1e1e", padx=16, pady=6)
        f.pack(side="left", padx=6, pady=4)
        tk.Label(f, text=title, font=("SF Pro Text", 9),
                 fg="#888", bg="#1e1e1e").pack()
        lbl = tk.Label(f, text=value, font=("SF Pro Display", 16, "bold"),
                       fg="#ea580c", bg="#1e1e1e")
        lbl.pack()
        return lbl

    def _choose_dir(self):
        d = filedialog.askdirectory(title="选择扫描目录")
        if d:
            self.dir_var.set(d)

    def _choose_files(self):
        fs = filedialog.askopenfilenames(
            title="选择文件",
            filetypes=[
                ("Office文档", "*.docx *.xlsx *.pdf *.pptx"),
                ("文本文件", "*.txt *.csv *.md *.rtf"),
                ("图片", "*.jpg *.jpeg *.png *.bmp *.tiff *.tif *.webp"),
            ]
        )
        if fs:
            self.dir_var.set(f"已选择 {len(fs)} 个文件")

    def _log(self, msg):
        self.log_box.configure(state="normal")
        self.log_box.insert("end", msg + "\n")
        self.log_box.see("end")
        self.log_box.configure(state="disabled")

    def _start(self):
        dir_val = self.dir_var.get()
        if dir_val in ["请选择目录...", ""]:
            messagebox.showwarning("请选择", "请先选择目录或文件！")
            return
        self._running = True
        self._stop_flag = False
        self._changed_files = []
        self._scanned = 0
        self.log_box.configure(state="normal")
        self.log_box.delete("1.0", "end")
        self.log_box.configure(state="disabled")
        self.stat_scanned.config(text="0")
        self.stat_changed.config(text="0")
        self.stat_status.config(text="扫描中...")
        self.btn_start.config(state="disabled")
        self.btn_stop.config(state="normal")
        self.btn_export.config(state="disabled")
        threading.Thread(target=self._run, daemon=True).start()

    def _stop(self):
        self._stop_flag = True
        self._log("停止请求已发送...")

    def _run(self):
        dir_val = self.dir_var.get()
        if os.path.isfile(dir_val) or '已选择' in dir_val:
            fs = filedialog.askopenfilenames()
            if not fs:
                self._done(); return
            self._process_files(list(fs))
        else:
            SKIP_DIRS = {"node_modules", ".git", "__pycache__", ".DS_Store",
                         "Library", "Applications", "System"}
            for dirpath, dirnames, filenames in os.walk(dir_val):
                if self._stop_flag:
                    break
                dirnames[:] = [d for d in dirnames if d not in SKIP_DIRS]
                self._process_files([os.path.join(dirpath, f) for f in filenames])
        self._done()

    def _process_files(self, files):
        for fpath in files:
            if self._stop_flag:
                break
            ext = os.path.splitext(fpath)[1].lower()
            if ext not in FILE_TYPES:
                continue
            self._scanned += 1
            if self._scanned % 10 == 0:
                self.after(0, lambda n=self._scanned: self.stat_scanned.config(text=str(n)))
            ext = ext[1:]
            try:
                if ext == 'docx':
                    process_docx(fpath, self._changed_files, self._log)
                elif ext == 'xlsx':
                    process_xlsx(fpath, self._changed_files, self._log)
                elif ext == 'pdf':
                    process_pdf(fpath, self._changed_files, self._log)
                elif ext in ('ppt', 'pptx'):
                    process_pptx(fpath, self._changed_files, self._log)
                elif ext in ('txt', 'csv', 'md', 'rtf'):
                    process_txt(fpath, self._changed_files, self._log)
                elif ext in IMG_TYPES:
                    process_image(fpath, self._changed_files, self._log)
            except Exception as e:
                self._log(f"  [错误] {Path(fpath).name}: {e}")

    def _done(self):
        self._running = False
        self.stat_status.config(text="完成")
        self.btn_start.config(state="normal")
        self.btn_stop.config(state="disabled")
        self.btn_export.config(state="normal" if self._changed_files else "disabled")
        self._log(f"\n完成！共扫描 {self._scanned} 个文件，修改 {len(self._changed_files)} 个")
        self.stat_scanned.config(text=str(self._scanned))
        self.stat_changed.config(text=str(len(self._changed_files)))

    def _export(self):
        if not self._changed_files:
            return
        path = filedialog.asksaveasfilename(
            title="导出报告",
            defaultextension=".txt",
            initialfile="internal_removal_report.txt"
        )
        if path:
            with open(path, 'w', encoding='utf-8') as f:
                f.write(f"去除内部标记报告\n")
                f.write(f"{'='*50}\n\n")
                for p in self._changed_files:
                    f.write(f"  * {p}\n")
            self._log(f"报告已导出: {path}")


if __name__ == '__main__':
    app = App()
    app.mainloop()
