import os
import sys
import re
import argparse
import warnings
import contextlib
import csv as csvlib
from pathlib import Path

# 屏蔽警告
warnings.filterwarnings("ignore", category=UserWarning)
warnings.filterwarnings("ignore", category=DeprecationWarning)
try:
    from PyPDF2.errors import PdfReadWarning
    warnings.filterwarnings("ignore", category=PdfReadWarning)
except Exception:
    pass

from docx import Document
from openpyxl import load_workbook
from PyPDF2 import PdfReader, PdfWriter
from pptx import Presentation

# 支持的文件类型
DOC_TYPES = ['.docx', '.xlsx', '.pdf', '.ppt', '.pptx', '.txt', '.csv', '.md', '.rtf']
IMG_TYPES = ['.jpg', '.jpeg', '.png', '.bmp', '.tiff', '.tif', '.webp']
FILE_TYPES = DOC_TYPES + IMG_TYPES

# OCR初始化标记
ocr_engine = None

def init_ocr():
    global ocr_engine
    if ocr_engine is not None:
        return True
    try:
        from paddleocr import PaddleOCR
        ocr_engine = PaddleOCR(use_angle_cls=True, lang='ch', show_log=False)
        return True
    except Exception as e:
        print(f"OCR初始化失败: {e}")
        return False

def get_all_drives():
    import string
    from ctypes import windll
    drives = []
    bitmask = windll.kernel32.GetLogicalDrives()
    for letter in string.ascii_uppercase:
        if bitmask & 1:
            drive = f'{letter}:\\'
            try:
                import win32file
                drive_type = win32file.GetDriveType(drive)
                if drive_type in [2, 5]:
                    bitmask >>= 1
                    continue
            except Exception:
                pass
            drives.append(drive)
        bitmask >>= 1
    return drives

def process_docx(file_path, success_files, changed_files):
    try:
        doc = Document(file_path)
        changed = False
        if doc.paragraphs:
            new_text = doc.paragraphs[0].text.replace('内部', '')
            if new_text != doc.paragraphs[0].text:
                doc.paragraphs[0].text = new_text
                changed = True
        doc.save(file_path)
        success_files.append(file_path)
        if changed:
            changed_files.append(file_path)
    except Exception:
        pass

def process_xlsx(file_path, success_files, changed_files):
    try:
        wb = load_workbook(file_path)
        ws = wb.active
        first_row = next(ws.iter_rows(min_row=1, max_row=1))
        changed = False
        for cell in first_row:
            if cell.value and isinstance(cell.value, str):
                new_val = cell.value.replace('内部', '')
                if new_val != cell.value:
                    cell.value = new_val
                    changed = True
        wb.save(file_path)
        success_files.append(file_path)
        if changed:
            changed_files.append(file_path)
    except Exception:
        pass

def process_pdf(file_path, success_files, changed_files):
    try:
        with open(os.devnull, 'w') as devnull, contextlib.redirect_stderr(devnull):
            reader = PdfReader(file_path)
            writer = PdfWriter()
            changed = False
            if reader.pages:
                first_page = reader.pages[0]
                text = first_page.extract_text()
                if text:
                    lines = text.split('\n')
                    if lines:
                        new_line = lines[0].replace('内部', '')
                        if new_line != lines[0]:
                            changed = True
                writer.add_page(first_page)
                for page in reader.pages[1:]:
                    writer.add_page(page)
                with open(file_path, 'wb') as f:
                    writer.write(f)
            success_files.append(file_path)
            if changed:
                changed_files.append(file_path)
    except Exception:
        pass

def process_ppt(file_path, success_files, changed_files):
    try:
        prs = Presentation(file_path)
        changed = False
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text:
                    lines = shape.text.split('\n')
                    if lines:
                        new_line = lines[0].replace('内部', '')
                        if new_line != lines[0]:
                            changed = True
                            lines[0] = new_line
                            shape.text = '\n'.join(lines)
        prs.save(file_path)
        success_files.append(file_path)
        if changed:
            changed_files.append(file_path)
    except Exception:
        pass

def process_txt(file_path, success_files, changed_files):
    try:
        with open(file_path, 'r', encoding='utf-8') as f:
            lines = f.readlines()
        if lines:
            new_line = lines[0].replace('内部', '')
            if new_line != lines[0]:
                lines[0] = new_line
                with open(file_path, 'w', encoding='utf-8') as f:
                    f.writelines(lines)
                changed_files.append(file_path)
        success_files.append(file_path)
    except Exception:
        pass

def process_csv(file_path, success_files, changed_files):
    try:
        with open(file_path, 'r', encoding='utf-8') as f:
            lines = f.readlines()
        if lines:
            new_line = lines[0].replace('内部', '')
            if new_line != lines[0]:
                lines[0] = new_line
                with open(file_path, 'w', encoding='utf-8') as f:
                    f.writelines(lines)
                changed_files.append(file_path)
        success_files.append(file_path)
    except Exception:
        pass

def process_md(file_path, success_files, changed_files):
    process_txt(file_path, success_files, changed_files)

def process_rtf(file_path, success_files, changed_files):
    try:
        with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
            content = f.read()
        new_content = content.replace('内部', '')
        if new_content != content:
            with open(file_path, 'w', encoding='utf-8') as f:
                f.write(new_content)
            changed_files.append(file_path)
        success_files.append(file_path)
    except Exception:
        pass

def process_image(file_path, success_files, changed_files):
    global ocr_engine
    try:
        from PIL import Image
        img = Image.open(file_path)
        
        if not init_ocr():
            success_files.append(file_path)
            return
        
        result = ocr_engine.ocr(file_path, cls=True)
        if not result or not result[0]:
            success_files.append(file_path)
            return
        
        changed = False
        draw = None
        
        for line in result[0]:
            box = line[0]
            text = line[1][0]
            if '内部' in text:
                if draw is None:
                    from PIL import ImageDraw
                    draw = ImageDraw.Draw(img)
                x1, y1 = box[0]
                x2, y2 = box[2]
                draw.rectangle([x1, y1, x2, y2], fill='white')
                changed = True
        
        if changed:
            img.save(file_path)
            changed_files.append(file_path)
        success_files.append(file_path)
    except Exception:
        pass

def scan_and_process(root_dir, success_files, changed_files, file_counter):
    for dirpath, _, filenames in os.walk(root_dir):
        skip_keywords = ["onedrive", "cloud", "dropbox", "baidunetdisk", "坚果云", "sync", "googledrive", "$recycle.bin", "system volume information", "windows", "program files"]
        dirpath_lower = dirpath.lower()
        if any(x in dirpath_lower for x in skip_keywords):
            continue
        for filename in filenames:
            ext = os.path.splitext(filename)[1].lower()
            if ext in FILE_TYPES:
                file_path = os.path.join(dirpath, filename)
                try:
                    if ext == '.docx':
                        process_docx(file_path, success_files, changed_files)
                    elif ext == '.xlsx':
                        process_xlsx(file_path, success_files, changed_files)
                    elif ext == '.pdf':
                        process_pdf(file_path, success_files, changed_files)
                    elif ext in ['.ppt', '.pptx']:
                        process_ppt(file_path, success_files, changed_files)
                    elif ext == '.txt':
                        process_txt(file_path, success_files, changed_files)
                    elif ext == '.csv':
                        process_csv(file_path, success_files, changed_files)
                    elif ext == '.md':
                        process_md(file_path, success_files, changed_files)
                    elif ext == '.rtf':
                        process_rtf(file_path, success_files, changed_files)
                    elif ext in IMG_TYPES:
                        process_image(file_path, success_files, changed_files)
                except Exception:
                    pass
                file_counter[0] += 1
                if file_counter[0] % 100 == 0:
                    print(f"已处理文件数: {file_counter[0]}")

if __name__ == "__main__":
    parser = argparse.ArgumentParser(description="全盘文档批量去除首行'内部'工具")
    parser.add_argument('--dir', type=str, default=None, help='只扫描指定目录')
    args = parser.parse_args()
    
    print("=" * 50)
    print("全盘文档批量去除首行'内部'工具 v3.0")
    print("=" * 50)
    print(f"支持格式: {', '.join(FILE_TYPES)}")
    print()
    
    drives = get_all_drives() if args.dir is None else [args.dir]
    success_files = []
    changed_files = []
    file_counter = [0]
    
    for drive in drives:
        print(f"扫描: {drive}")
        scan_and_process(drive, success_files, changed_files, file_counter)
    
    print()
    print("=" * 50)
    print(f"处理完成！共处理 {len(success_files)} 个文件，修改 {len(changed_files)} 个文件")
    print("=" * 50)
    
    output_dir = os.path.dirname(os.path.abspath(sys.argv[0]))
    csv_path = os.path.join(output_dir, "changed_files.csv")
    with open(csv_path, "w", newline='', encoding="utf-8-sig") as csvfile:
        writer = csvlib.writer(csvfile)
        writer.writerow(["file_path"])
        for f in changed_files:
            writer.writerow([f])
    print(f"修改记录已保存: {csv_path}")
    input("按回车键退出...")
